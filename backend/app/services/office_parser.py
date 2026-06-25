import io

WORDS_PER_PAGE = 400


def extract_docx(file_bytes: bytes) -> list[dict]:
    import docx
    doc = docx.Document(io.BytesIO(file_bytes))

    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                paragraphs.append(row_text)

    if not paragraphs:
        return []

    pages = []
    current_words = []
    current_word_count = 0
    page_num = 1

    for para in paragraphs:
        words = para.split()
        current_words.append(para)
        current_word_count += len(words)
        if current_word_count >= WORDS_PER_PAGE:
            pages.append({"page_number": page_num, "text": "\n\n".join(current_words)})
            page_num += 1
            current_words = []
            current_word_count = 0

    if current_words:
        pages.append({"page_number": page_num, "text": "\n\n".join(current_words)})

    return pages


def extract_xlsx(file_bytes: bytes) -> list[dict]:
    import openpyxl
    wb = openpyxl.load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)

    pages = []
    for page_num, sheet_name in enumerate(wb.sheetnames, start=1):
        ws = wb[sheet_name]
        rows_text = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            row_str = " | ".join(cells).strip(" |")
            if row_str:
                rows_text.append(row_str)

        if rows_text:
            pages.append({
                "page_number": page_num,
                "text": f"[Sheet: {sheet_name}]\n" + "\n".join(rows_text),
            })

    wb.close()
    return pages


def extract_pptx(file_bytes: bytes) -> list[dict]:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(file_bytes))

    pages = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        if texts:
            pages.append({
                "page_number": slide_num,
                "text": "\n\n".join(texts),
            })

    return pages

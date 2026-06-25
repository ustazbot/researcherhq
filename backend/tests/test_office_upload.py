import sys, os
sys.path.insert(0, os.path.join(os.path.dirname(__file__), '..'))

import io
import sqlite3
import pytest
from unittest.mock import patch, AsyncMock
from fastapi.testclient import TestClient
from app.services.auth_service import create_jwt


def make_headers(user_id="user-1", email="u1@test.com"):
    token = create_jwt({"user_id": user_id, "email": email})
    return {"Authorization": f"Bearer {token}"}


def make_docx_bytes(text="Hello world " * 50) -> bytes:
    import docx
    doc = docx.Document()
    doc.add_paragraph(text)
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def make_xlsx_bytes() -> bytes:
    import openpyxl
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    ws.append(["Name", "Score", "Department", "Notes", "Grade"])
    # 30 rows × 5 columns of words → well over MIN_CHUNK_SIZE (100 words)
    for i in range(30):
        ws.append([f"Student{i}", str(80 + i % 20), "Engineering", f"note{i} detail about performance", "A"])
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def make_pptx_bytes(slide_text="Slide one content") -> bytes:
    from pptx import Presentation
    prs = Presentation()
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = slide_text
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


DOCX_MIME = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
XLSX_MIME = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"


@pytest.fixture
def client_with_project(tmp_path):
    db_path = str(tmp_path / "test.db")
    with patch("app.database._db_path", db_path):
        from app.database import init_db
        init_db(db_path)
        from app.main import app
        with TestClient(app) as c:
            headers = make_headers()
            proj_r = c.post("/projects", json={"title": "Test", "research_mode": "general"}, headers=headers)
            project_id = proj_r.json()["id"]
            yield c, project_id, headers


@pytest.fixture
def pro_client_with_project(tmp_path):
    db_path = str(tmp_path / "test.db")
    with patch("app.database._db_path", db_path):
        from app.database import init_db
        init_db(db_path)
        from app.main import app
        with TestClient(app) as c:
            headers = make_headers(user_id="pro-user", email="pro@test.com")
            conn = sqlite3.connect(db_path)
            conn.execute(
                "INSERT OR IGNORE INTO users (id, email, password_hash, tier) VALUES ('pro-user', 'pro@test.com', 'x', 'pro')"
            )
            conn.commit()
            conn.close()
            proj_r = c.post("/projects", json={"title": "Pro Project", "research_mode": "general"}, headers=headers)
            project_id = proj_r.json()["id"]
            yield c, project_id, headers


# --- Test 1: DOCX upload ---
def test_docx_upload_success(client_with_project):
    with patch("app.routers.documents._embed_and_store_chunks", new_callable=AsyncMock):
        client, project_id, headers = client_with_project
        r = client.post(
            "/documents/upload-office",
            data={"project_id": project_id, "category": "artikel"},
            files={"file": ("test.docx", make_docx_bytes(), DOCX_MIME)},
            headers=headers,
        )
    assert r.status_code == 201
    data = r.json()
    assert data["file_type"] == "docx"
    assert data["chunk_count"] >= 1
    assert data["status"] == "uploaded"


# --- Test 2: XLSX upload ---
def test_xlsx_upload_success(client_with_project):
    with patch("app.routers.documents._embed_and_store_chunks", new_callable=AsyncMock):
        client, project_id, headers = client_with_project
        r = client.post(
            "/documents/upload-office",
            data={"project_id": project_id, "category": "data"},
            files={"file": ("data.xlsx", make_xlsx_bytes(), XLSX_MIME)},
            headers=headers,
        )
    assert r.status_code == 201
    data = r.json()
    assert data["file_type"] == "xlsx"
    assert data["chunk_count"] >= 1


# --- Test 3: PPTX upload ---
def test_pptx_upload_success(client_with_project):
    with patch("app.routers.documents._embed_and_store_chunks", new_callable=AsyncMock):
        client, project_id, headers = client_with_project
        r = client.post(
            "/documents/upload-office",
            data={"project_id": project_id, "category": "catatan_sv"},
            files={"file": ("slides.pptx", make_pptx_bytes(), PPTX_MIME)},
            headers=headers,
        )
    assert r.status_code == 201
    data = r.json()
    assert data["file_type"] == "pptx"


# --- Test 4: Invalid MIME type ---
def test_invalid_mime_type(client_with_project):
    client, project_id, headers = client_with_project
    r = client.post(
        "/documents/upload-office",
        data={"project_id": project_id, "category": "artikel"},
        files={"file": ("evil.txt", b"hello", "text/plain")},
        headers=headers,
    )
    assert r.status_code == 400
    assert "DOCX, XLSX, and PPTX" in r.json()["detail"]


# --- Test 5: File too large ---
def test_file_too_large(client_with_project):
    client, project_id, headers = client_with_project
    big_bytes = b"x" * (20 * 1024 * 1024 + 1)
    r = client.post(
        "/documents/upload-office",
        data={"project_id": project_id, "category": "data"},
        files={"file": ("big.docx", big_bytes, DOCX_MIME)},
        headers=headers,
    )
    assert r.status_code == 413
    assert "20MB" in r.json()["detail"]


# --- Test 6: Empty file / no text content ---
def test_empty_file_no_text(client_with_project):
    import docx
    doc = docx.Document()
    buf = io.BytesIO()
    doc.save(buf)
    empty_bytes = buf.getvalue()

    client, project_id, headers = client_with_project
    r = client.post(
        "/documents/upload-office",
        data={"project_id": project_id, "category": "artikel"},
        files={"file": ("empty.docx", empty_bytes, DOCX_MIME)},
        headers=headers,
    )
    assert r.status_code == 422
    assert "No text content" in r.json()["detail"]


# --- Test 7: Invalid category ---
def test_invalid_category(client_with_project):
    client, project_id, headers = client_with_project
    r = client.post(
        "/documents/upload-office",
        data={"project_id": project_id, "category": "bogus"},
        files={"file": ("test.docx", make_docx_bytes(), DOCX_MIME)},
        headers=headers,
    )
    assert r.status_code == 400
    assert "Invalid category" in r.json()["detail"]


# --- Test 8: Project not found / wrong user ---
def test_project_not_found(client_with_project):
    client, _, headers = client_with_project
    r = client.post(
        "/documents/upload-office",
        data={"project_id": "nonexistent-id", "category": "artikel"},
        files={"file": ("test.docx", make_docx_bytes(), DOCX_MIME)},
        headers=headers,
    )
    assert r.status_code == 404
    assert "not found" in r.json()["detail"].lower()


# --- Test 9: Free tier doc limit ---
def test_free_tier_doc_limit(client_with_project):
    with patch("app.routers.documents._embed_and_store_chunks", new_callable=AsyncMock):
        client, project_id, headers = client_with_project
        r1 = client.post(
            "/documents/upload-office",
            data={"project_id": project_id, "category": "artikel"},
            files={"file": ("doc1.docx", make_docx_bytes(), DOCX_MIME)},
            headers=headers,
        )
        assert r1.status_code == 201

        r2 = client.post(
            "/documents/upload-office",
            data={"project_id": project_id, "category": "artikel"},
            files={"file": ("doc2.docx", make_docx_bytes(), DOCX_MIME)},
            headers=headers,
        )
    assert r2.status_code == 403
    assert "limit" in r2.json()["detail"].lower()


# --- Test 10: panduan_format category accepted ---
def test_panduan_format_category(client_with_project):
    with patch("app.routers.documents._embed_and_store_chunks", new_callable=AsyncMock):
        client, project_id, headers = client_with_project
        r = client.post(
            "/documents/upload-office",
            data={"project_id": project_id, "category": "panduan_format"},
            files={"file": ("guide.docx", make_docx_bytes(), DOCX_MIME)},
            headers=headers,
        )
    assert r.status_code == 201
    data = r.json()
    assert data["status"] == "uploaded"
